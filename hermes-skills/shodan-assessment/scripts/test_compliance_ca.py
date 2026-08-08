#!/usr/bin/env python3
"""test_compliance_ca.py — the Canadian regime set, and the claims it must NEVER make.

WHY THIS EXISTS. `CA_COMPLIANCE_REFERENCE.md` §6.4 carries an explicit list of statements that are
WRONG in a Canadian deck. Every one of them is the kind of error a reader who works at a bank spots
instantly and cannot unsee — an OSFI "fine", a live CCSPA deadline, a PIPEDA penalty attached to
the breach itself. Research produced that list; this turns it into a build gate, because a rule
that lives only in a markdown file is a rule that gets broken by the next edit.

It also pins the structural half: the regime set FOLLOWS THE JURISDICTION through one registry,
the way build_findings_deck.js already selects its framework set from the estate's country. Canada
used to fall through to the generic ISO 27001 / NIST CSF default.

    python test_compliance_ca.py          # exits non-zero on any failure
"""
import importlib.util
import json
import os
import re
import subprocess
import sys
import tempfile

HERE = os.path.dirname(os.path.abspath(__file__))
FAILED = []


def _load(name, filename):
    spec = importlib.util.spec_from_file_location(name, os.path.join(HERE, filename))
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)
    return mod


def check(cond, msg):
    print(("  PASS  " if cond else "  FAIL  ") + msg)
    if not cond:
        FAILED.append(msg)


CE = _load("compliance_enrich", "compliance_enrich.py")
CC = _load("compliance_clarify", "compliance_clarify.py")

print("=" * 78)
print("  Canadian compliance — regime set + the claims a Canadian deck must never make")
print("=" * 78)

# ---------------------------------------------------------------- §1 the registry ------------
print("\n[1] the regime set follows the JURISDICTION, through one registry")
check(CE.regimes_for("CA") and CE.regimes_for("CA") != CE.regimes_for("EU"),
      "CA resolves to its own regime set, not the EU one")
check(all(k in CE.FIXED for k in CE.regimes_for("CA")),
      "every CA regime has FIXED regulatory facts (no empty deck sections)")
check(set(CE.decks_for("CA")) <= set(CE.regimes_for("CA")),
      "every CA deck key is a graded regime")
check(CE.jurisdiction("ZZ")[0] == CE.DEFAULT_JURISDICTION,
      "an unknown jurisdiction falls back to a real regime set, never an empty one")
check(CE.regimes_for("ca") == CE.regimes_for("CA"), "lookup is case-insensitive")

CA, _ = CE.build("Royal Bank of Canada", "en", None, "CA")
EU, _ = CE.build("Acme GmbH", "en", None, "EU")
blob = json.dumps(CA, ensure_ascii=False)

check(CA["jurisdiction"] == "CA" and CA["order"] == CE.regimes_for("CA"),
      "compliance.json carries its own order, so the deck builders need no constant")
check("CA_COMPLIANCE_REFERENCE.md" in CA["source"],
      "the Canadian output cites the Canadian reference")
check(EU["order"] == ["nis2", "cra", "aiact"],
      "REGRESSION: the EU set is unchanged")

# ------------------------------------------------- §2 the reference's §6.4 forbidden list -----
print("\n[2] §6.4 — what must NEVER appear in a Canadian deck")

# Assert the PROPERTY on EVERY OSFI regime, not a regex hunting for a symptom. The first version
# anchored on the literal "osfi" with an 80-character window, which never reached the penalty field
# 400 characters later — a negative test injecting "$5,000,000 fine" sailed straight through it.
osfi_keys = [k for k in CA["order"] if k.startswith("osfi")]
check(len(osfi_keys) >= 4, "the OSFI block is present (%d guidelines)" % len(osfi_keys))
bad_osfi = [k for k in osfi_keys
            if str((CA["regimes"][k]["penalty"] or {}).get("essential", "")).strip()
            != "No monetary penalty"
            or re.search(r"\$[\d,]+", json.dumps(CA["regimes"][k]["penalty"], ensure_ascii=False))]
check(not bad_osfi,
      "no OSFI fine anywhere — OSFI's tools are supervisory, not monetary%s"
      % ((" [offenders: %s]" % ", ".join(bad_osfi)) if bad_osfi else ""))

ccspa = json.dumps(CA["regimes"]["ccspa"], ensure_ascii=False)
check("NOT IN FORCE" in ccspa.upper(), "the CCSPA is labelled NOT IN FORCE")
check("1,000,000" not in ccspa and "$1M" not in ccspa,
      "no $1,000,000 individual CCSPA AMP (halved to $500,000 at committee, 26 Feb 2026)")
check(re.search(r"ceiling", ccspa, re.I) is not None,
      "CCSPA penalties are described as ceilings on future regulations, not a live tariff")
check("72" not in ccspa or re.search(r"not to exceed 72|planning only|prescribed by", ccspa, re.I),
      "no live 72-hour CCSPA clock — it is a future ceiling, stated as planning only")
_cc = CA["regimes"]["ccspa"]
check(_cc["classification"] != "Applies" and _cc["applies"] is not True
      and "NOT IN FORCE" in (str(_cc["instrument"]) + str(_cc["regulates"])).upper()
      and not re.search(r"\bApplies\b", str(_cc["regulates"])),
      "the CCSPA is never presented as applying today (classification, applies flag, "
      "instrument and scope line all agree)")

pipeda = json.dumps(CA["regimes"]["pipeda"], ensure_ascii=False)
check("knowingly" in pipeda.lower() and "no penalty for the breach itself" in pipeda.lower(),
      "the PIPEDA fine is tied to KNOWINGLY contravening s.10.1/10.3 — not to the breach or cl. 4.7")

law25 = json.dumps(CA["regimes"]["law25"], ensure_ascii=False)
check(re.search(r"open constitutional question", law25, re.I) is not None,
      "Law 25 applicability to a federally chartered bank is stated as an OPEN question")
check(not re.search(r"law 25 (applies to|binds) (a )?(federal|frfi|bank)", blob, re.I),
      "nothing asserts that Law 25 binds a federally chartered bank")

check(not re.search(r"itsg-33[^.]{0,40}(mandatory|required)", blob, re.I),
      "ITSG-33 is never called mandatory")
check(not re.search(r"baseline (cyber security )?controls", blob, re.I),
      "CCCS Baseline Controls (an SME framework) are not offered to a bank")
check("C-36" not in blob, "no Bill C-36 figure")

# ----------------------------------------------------------- §3 the facts that MUST appear ----
print("\n[3] the live facts a Canadian FRFI is actually buying")
e21 = json.dumps(CA["regimes"]["osfi_e21"], ensure_ascii=False)
check("2026-09-01" in e21, "E-21 full adherence on 1 Sep 2026 — the nearest hard deadline")
check("2027-09-01" in e21, "E-21 scenario testing complete by 1 Sep 2027")
check("2025-03-31" in json.dumps(CA["regimes"]["osfi_b10"], ensure_ascii=False),
      "B-10 foreign-branch date (31 Mar 2025) modelled separately from 1 May 2024")
check("24 hour" in json.dumps(CA["regimes"]["osfi_incident"]).lower()
      or "24 hours" in json.dumps(CA["regimes"]["osfi_incident"]).lower(),
      "the OSFI 24-hour incident clock — the tightest in the Canadian set")
check("minimizing its attack surface" in json.dumps(CA["regimes"]["osfi_b13"]),
      "B-13 §3.2.4 quoted verbatim — the clause this product maps onto")
check(str(CA.get("disclaimer", "")).lower().startswith("not legal advice"),
      "the output carries the not-legal-advice disclaimer")

# ------------------------------------------------------------------ §4 vendor neutrality ------
print("\n[4] vendor neutrality — the fallback path the brand gate never rendered")
check("Colt" not in blob, "no carrier brand in the Canadian deterministic output")
check("Colt" not in json.dumps(EU, ensure_ascii=False),
      "REGRESSION: no carrier brand in the EU deterministic output either")
check("Colt" not in CE.PROMPT,
      "no carrier brand in the PROMPT — it is an instruction that reaches the reader via the model")
check("colt" in CA["regimes"]["osfi_b13"],
      "the `colt` JSON key is NOT renamed — it is a lookup key the deck builder reads")

# --------------------------------------------------------------------- §5 clarify -------------
print("\n[5] clarification questions follow the jurisdiction")
qca = CC.build(CA)
qeu = CC.build(EU)
ids_ca = [q["id"] for q in qca["questions"]]
ids_eu = [q["id"] for q in qeu["questions"]]
check("frfi" in ids_ca and "quebec" in ids_ca, "CA asks FRFI status and Quebec presence")
check("frfi" not in ids_eu and "quebec" not in ids_eu,
      "an EU customer is never asked whether it is a Canadian FRFI")
check(all(q.get("maps_to") for q in qca["questions"]),
      "every question is machine-actionable (carries maps_to)")
check(set(CA["order"]) <= set(qca["summary"]),
      "the clarify summary covers every graded regime")

# ------------------------------------------------------------------ §6 the decks build --------
print("\n[6] every advertised deck actually renders")
tmp = tempfile.mkdtemp(prefix="ca_compliance_")
cpath = os.path.join(tmp, "compliance.json")
json.dump(CA, open(cpath, "w", encoding="utf-8"), ensure_ascii=False)
built = 0
for key in list(CA["decks"]) + ["roadmap"]:
    outp = os.path.join(tmp, "%s.pptx" % key)
    r = subprocess.run(["node", os.path.join(HERE, "build_compliance_deck.js"), cpath, outp, key],
                       capture_output=True, text=True)
    ok = r.returncode == 0 and os.path.exists(outp) and os.path.getsize(outp) > 40000
    if not ok:
        print("        %s: %s" % (key, (r.stderr or "").strip()[:160]))
    built += 1 if ok else 0
check(built == len(CA["decks"]) + 1,
      "all %d Canadian decks render (%d built)" % (len(CA["decks"]) + 1, built))

# The TITLE SLIDE is a different code path from the content slides — fixing the eyebrow in
# content() did NOT fix it, and the roadmap cover read "EU DIGITAL & CYBER COMPLIANCE" over
# "NIS2 · CRA · EU AI Act" on a Canadian bank's deck. Read the rendered text, not the source.
try:
    from pptx import Presentation as _P
    _rd = os.path.join(tmp, "roadmap.pptx")
    _txt = " ".join(r.text for _s in _P(_rd).slides for _sh in _s.shapes
                    if _sh.has_text_frame for _p in _sh.text_frame.paragraphs for r in _p.runs)
    _leak = [n for n in ("EU DIGITAL", "NIS2", "Cyber Resilience Act", "EU AI Act",
                         "EU primary law") if n in _txt]
    check(not _leak, "no EU chrome on the Canadian roadmap deck%s"
          % ((" [found: %s]" % ", ".join(_leak)) if _leak else ""))
    check("CANADIAN DIGITAL" in _txt and "OSFI" in _txt,
          "the Canadian roadmap cover names Canada and OSFI")
except ImportError:
    print("        (python-pptx unavailable - deck text check skipped)")

html = os.path.join(tmp, "report.html")
r = subprocess.run(["node", os.path.join(HERE, "build_compliance_html.js"), cpath, html],
                   capture_output=True, text=True)
ok_html = r.returncode == 0 and os.path.exists(html)
check(ok_html, "the animated compliance report renders for Canada")
if ok_html:
    h = open(html, encoding="utf-8").read()
    # THE CHECK THAT WAS MISSING. The old version asserted the report RENDERED and carried no
    # undefined/Colt — never that it named the right regimes. So build_compliance_html.js kept its
    # four hardcoded ["nis2","cra","aiact"] arrays and its "EU Compliance" title, and a Canadian
    # bank received an EU report. A check that cannot see the thing it checks is not a check.
    eu_leak = [n for n in ("NIS2", "Cyber Resilience Act", "EU AI Act", "EU Compliance") if n in h]
    check(not eu_leak,
          "the Canadian report names NO EU regime%s"
          % ((" [found: %s]" % ", ".join(eu_leak)) if eu_leak else ""))
    check(all(n in h for n in ("OSFI", "PIPEDA", "CCSPA")),
          "the Canadian report names OSFI, PIPEDA and the CCSPA")
    check("Canadian Compliance" in h, "the report TITLE follows the jurisdiction")
    check(not re.search(r"undefined|NaN|\[object Object\]", h),
          "no undefined/NaN/[object Object] leaks into the Canadian report")
    check("Colt" not in h, "no carrier brand in the rendered Canadian report")

# ------------------------------------------------- §7 findings deck framework set --------------
print("\n[7] the findings deck cites Canadian regulators for a Canadian estate")
fd = open(os.path.join(HERE, "build_findings_deck.js"), encoding="utf-8").read()
ca_block_raw = fd[fd.find('cc === "CA"'):fd.find('EU.indexOf(cc)')]
# Strip the comments FIRST. The block's own comment explains why the CCSPA is excluded, so a naive
# grep finds "CCSPA" in the justification and reports the defect it was written to prevent — a
# check that reads its own explanation as evidence. The brand gate strips comments for this reason.
ca_block = re.sub(r"//[^\n]*", "", ca_block_raw)
check('cc === "CA"' in fd, "the findings deck has a CA branch")
check("OSFI B-13" in ca_block and "PIPEDA" in ca_block,
      "a Canadian estate is mapped to OSFI B-13 and PIPEDA")
check("CCSPA" not in ca_block,
      "the CCSPA is NOT offered as an applicable framework — Part 2 is not in force")
check("NIS2" not in ca_block and "GDPR" not in ca_block,
      "no EU regime is cited at a Canadian bank")

print("\n" + "=" * 78)
if FAILED:
    print("  FAILED (%d):" % len(FAILED))
    for f in FAILED:
        print("    - " + f)
    sys.exit(1)
print("  ALL CHECKS PASSED — the Canadian set is correct and cannot silently regress")
print("=" * 78)
