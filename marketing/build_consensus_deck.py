#!/usr/bin/env python3
"""build_consensus_deck.py — the 4-model consensus deep-dive, in the S4biz template.

WHY A SCRIPT AND NOT A HAND-MADE FILE: standing rule 6 — every deliverable is a re-runnable
script committed to the repo. Re-run it and the deck rebuilds identically.

The template is NOT re-implemented: every coordinate, colour and font below was READ off
`S4biz_Sovereign_Cyber_Cloud_Capability_Brief.pptx` (see extract notes in the palette block), and
the two background images are lifted from that file's media parts. So the output is visually the
same deck family, not an approximation of it.

    python marketing/build_consensus_deck.py [--out PATH] [--template PATH]

EVERY factual claim on the science slides is sourced. Citations were verified against the primary
papers (arXiv/proceedings), including two corrections to my own initial premises — see SOURCES.
"""
import argparse
import io
import os
import zipfile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------------------------
# PALETTE + TYPE — read directly off the template, not invented.
#   fills   14161F (card) · 1B1F2C (panel) · 22D3EE (cyan) · 8B5CF6 (violet) · 4F46E5 (indigo)
#   line    2B3042      text  FFFFFF / C7CDDA (body) / 8E97A8 (muted)
#   fonts   Arial Black (display) · Arial (body) · Consolas (mono, eyebrow + footer)
# ---------------------------------------------------------------------------------------------
INK      = RGBColor(0x14, 0x16, 0x1F)
PANEL    = RGBColor(0x1B, 0x1F, 0x2C)
LINE     = RGBColor(0x2B, 0x30, 0x42)
CYAN     = RGBColor(0x22, 0xD3, 0xEE)
VIOLET   = RGBColor(0x8B, 0x5C, 0xF6)
INDIGO   = RGBColor(0x4F, 0x46, 0xE5)
RED      = RGBColor(0xFF, 0x4D, 0x6D)
AMBER    = RGBColor(0xFF, 0x9F, 0x1C)
GREEN    = RGBColor(0x3D, 0xD6, 0x8C)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
BODY     = RGBColor(0xC7, 0xCD, 0xDA)
MUTED    = RGBColor(0x8E, 0x97, 0xA8)

DISPLAY, TEXT, MONO = "Arial Black", "Arial", "Consolas"
W, H = 13.333, 7.5


def _tb(slide, x, y, w, h, text, size=11, color=BODY, font=TEXT, bold=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=1.0, wrap=True):
    """Text box with zero internal padding (the skill's rule: margin 0 or nothing aligns)."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = space
        r = p.add_run()
        r.text = line
        r.font.size, r.font.name, r.font.bold = Pt(size), font, bold
        r.font.color.rgb = color
    return box


def _rect(slide, x, y, w, h, fill=INK, line=LINE, lw=1.0, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    s.shadow.inherit = False
    if s.has_text_frame:
        s.text_frame.word_wrap = True
    return s


class Deck:
    def __init__(self, template):
        self.prs = Presentation()
        self.prs.slide_width, self.prs.slide_height = Inches(W), Inches(H)
        with zipfile.ZipFile(template) as z:
            self.bg_title = z.read("ppt/media/Slide-1-image-1.png")
            self.bg_body = z.read("ppt/media/Slide-3-image-1.png")
        self.n = 0

    def slide(self, eyebrow, title, title_tail=None, sub=None, footer="", hero=False):
        """One slide carrying the template's exact chrome. Returns the slide.

        ARGUMENT ORDER IS (eyebrow, title, title_tail, sub, footer) — it reads the way the slide
        reads, top to bottom. The first version declared sub and footer BEFORE title_tail while
        every call site passed them in visual order, so the footer rendered as the sub-heading and
        the sub rendered as the violet tail of the title. It cost a full render cycle to see.
        Same class as every other defect in this deck: assuming a helper's contract instead of
        reading it.
        """
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        # background FIRST so everything else paints over it
        s.shapes.add_picture(io.BytesIO(self.bg_title if hero else self.bg_body),
                             0, 0, Inches(W), Inches(H))
        self.n += 1

        # --- wordmark + the three squares -----------------------------------------------------
        # The template's own coordinates put "S4BIZ" at x11.33 and the squares at x11.91, and at
        # 18pt Arial Black the wordmark is ~0.7in wide — so the squares land on top of "BIZ".
        # Right-align the wordmark to end at 11.85 instead, which keeps both elements and their
        # sizes and removes the collision. Measured, not eyeballed.
        wm = _tb(s, 10.20, 0.28, 1.50, 0.45, "", 18, WHITE, DISPLAY, True, PP_ALIGN.RIGHT)
        p = wm.text_frame.paragraphs[0]
        for txt, col in (("S4", VIOLET), ("BIZ", WHITE)):
            r = p.add_run()
            r.text = txt
            r.font.size, r.font.name, r.font.bold = Pt(18), DISPLAY, True
            r.font.color.rgb = col
        for i, col in enumerate((INDIGO, VIOLET, CYAN)):
            _rect(s, 11.91 + i * 0.26, 0.34, 0.34, 0.34, col, None)

        _tb(s, 0.55, 0.50, 11.00, 0.30, eyebrow.upper(), 10.5, CYAN, MONO, True)

        if hero:
            tb = _tb(s, 0.50, 1.75, 12.30, 2.55, "", 38, WHITE, DISPLAY, True)
            tf = tb.text_frame
            for i, (txt, col) in enumerate(title):
                para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                para.line_spacing = 0.95
                r = para.add_run()
                r.text = txt
                r.font.size, r.font.name, r.font.bold = Pt(38), DISPLAY, True
                r.font.color.rgb = col
        else:
            tb = _tb(s, 0.50, 0.82, 12.30, 0.70, "", 30, WHITE, DISPLAY, True)
            para = tb.text_frame.paragraphs[0]
            r = para.add_run()
            r.text = title
            r.font.size, r.font.name, r.font.bold = Pt(30), DISPLAY, True
            r.font.color.rgb = WHITE
            if title_tail:
                r2 = para.add_run()
                r2.text = title_tail
                r2.font.size, r2.font.name, r2.font.bold = Pt(30), DISPLAY, True
                r2.font.color.rgb = VIOLET
            if sub:
                _tb(s, 0.57, 1.55, 12.10, 0.35, sub, 11.5, BODY, TEXT)

        _tb(s, 0.55, 7.04, 10.50, 0.30, footer.upper(), 8.5, MUTED, MONO)
        _tb(s, 11.93, 7.04, 0.85, 0.30, ">> %02d" % self.n, 8.5, CYAN, MONO)
        return s

    def save(self, path):
        self.prs.save(path)
        return path


def card(s, x, y, w, h, kicker, kcol, head, body, hcol=WHITE, bsize=10, fill=INK):
    """The template's card: panel + a coloured kicker line + head + body.

    THE HEAD IS A FIXED BOX AND IT IS ARITHMETIC, not taste. The head sits at y+0.50 and the body
    starts at y+0.95, so the head has 0.45in. At 14pt Arial Black a line is about 0.233in: one line
    fits comfortably, two just fit, three overlap the body. A three-line head shipped in the first
    render of the cybergod pitch (the MITRE / Diamond / Admiralty card) and printed straight
    through the paragraph beneath it.

    Same defect class as the site header row and the deck title row, both of which this repository
    has already paid for. Refuse it at build time rather than hope somebody looks at the render.
    """
    if head.count("\n") + 1 > 2:
        raise SystemExit("[X] card head is %d lines and overlaps the body (max 2): %r"
                         % (head.count("\n") + 1, head))
    # A TWO-LINE HEAD MOVES THE BODY DOWN. At 14pt a line is ~0.233in, so two lines need 0.47in
    # while the gap to the body is 0.45: the second line and the first line of body text touch.
    # Measured on the render rather than assumed. A one-line head is unaffected, so the two decks
    # that only use one-line heads are byte-identical to before this change.
    drop = 0.24 if head.count("\n") else 0.0
    _rect(s, x, y, w, h, fill, LINE)
    _tb(s, x + 0.26, y + 0.22, w - 0.52, 0.24, kicker.upper(), 9, kcol, MONO, True)
    _tb(s, x + 0.26, y + 0.50, w - 0.52, 0.38 + drop, head, 14, hcol, DISPLAY, True)
    _tb(s, x + 0.26, y + 0.95 + drop, w - 0.52, h - 1.15 - drop, body, bsize, BODY, TEXT,
        space=1.28)


def bullets(s, x, y, w, items, gap=0.52, dot=CYAN, size=10.5, dotsize=0.11):
    """Bullet list whose spacing FOLLOWS THE WRAP.

    A fixed per-item offset makes a two-line bullet crowd the one beneath it: the QA render
    showed 57px between single-line bullets and 33px after a wrapped one, so the lower bullets
    visually clumped into a block. Estimate the line count from the measured character capacity
    and advance by that instead — same arithmetic discipline as the header row.
    """
    cpl = max(20, int((w - 0.30) / (size * 0.0064)))   # calibrated on the 110dpi render
    yy = y
    for it in items:
        lines = max(1, -(-len(it) // cpl))
        _rect(s, x, yy + 0.055, dotsize, dotsize, dot, None)
        _tb(s, x + 0.30, yy, w - 0.30, lines * 0.20 + 0.04, it, size, BODY, TEXT, space=1.22)
        yy += max(gap, lines * 0.205 + 0.155)


def stat(s, x, y, w, value, label, col=CYAN, vsize=30):
    _tb(s, x, y, w, 0.55, value, vsize, col, DISPLAY, True)
    _tb(s, x, y + 0.58, w, 0.52, label, 9, MUTED, MONO, space=1.2)


def build(template, out):
    d = Deck(template)

    # =========================================================================================
    # 01 — TITLE
    # =========================================================================================
    s = d.slide("S4biz · Cybergod · engineering deep dive · 7 Aug 2026",
                [("FOUR MODELS", WHITE), ("CHECK THE WORK.", WHITE), ("CODE DECIDES.", VIOLET)],
                footer="Stars4business OÜ · Cybergod LLC · engineering note · confidential",
                hero=True)
    _tb(s, 0.55, 1.35, 9.00, 0.32, "> ship.py --stage --panel 4 --gate deterministic",
        12, VIOLET, MONO)
    _tb(s, 0.57, 4.45, 11.50, 0.45,
        "How the cybergod.ai release gate works · what the evidence actually says · "
        "and what we do not claim.", 12, BODY, TEXT)
    for i, (t, sub2, col) in enumerate([
            ("2", "soldiers\nbuild the case", CYAN),
            ("2", "auditors\nattack the case", VIOLET),
            ("4", "vendors\nno shared failure", INDIGO),
            ("35", "checks decide\nmodels only advise", GREEN),
            ("195s", "one command\ntest → stage → prod", AMBER)]):
        x = 0.55 + i * 2.49
        _rect(s, x, 5.45, 2.29, 1.35, INK, LINE)
        _rect(s, x, 5.45, 2.29, 0.06, col, None)
        _tb(s, x + 0.16, 5.61, 1.99, 0.50, t, 22, col, DISPLAY, True)
        _tb(s, x + 0.16, 6.07, 1.99, 0.66, sub2, 9.5, WHITE, TEXT, True, space=1.15)

    # =========================================================================================
    # 02 — WHAT SHIPPED TODAY
    # =========================================================================================
    s = d.slide("Today's release · three defects, found by a customer",
                "WHAT SHIPPED", " TODAY.",
                "Three visible defects. Each one is now a build gate, so it cannot come back.",
                "release 1df00ba · good-20260807-223843 · confidential")
    card(s, 0.55, 2.20, 3.95, 2.35, "defect 01", RED, "Gibberish text",
         "The live page read “somebody&rsquo;s inbox”. JSX parses entities in literal "
         "text; a string routed through the translator is a JS string, so React escapes it and "
         "the raw code reaches the screen.")
    card(s, 4.69, 2.20, 3.95, 2.35, "defect 02", RED, "A block, twice",
         "The three FAQ cards rendered two times — once from the keyed dictionary, once as "
         "raw English literals. Same content, two key spaces. Visible to every visitor.")
    card(s, 8.83, 2.20, 3.95, 2.35, "defect 03", RED, "Dead menu on Android",
         "Eight taps on the phone, no menu, ever — while the identical code worked perfectly "
         "on desktop. Root cause on the next slide; it was not a touch problem.")
    _rect(s, 0.55, 4.85, 12.23, 1.95, PANEL, LINE)
    _tb(s, 0.85, 5.05, 11.60, 0.30, "AND THE PART THAT MATTERS", 10.5, CYAN, MONO, True)
    _tb(s, 0.85, 5.40, 11.60, 1.20,
        "Each fix shipped with a gate that fails the build if the defect returns. The new "
        "entity gate then found 20 MORE live instances nobody had reported — including "
        "“Swipe the map sideways to explore &rarr;”, in all six languages.\n"
        "A defect you fix once comes back. A defect you encode as a rule cannot.",
        11.5, BODY, TEXT, space=1.35)

    # =========================================================================================
    # 03 — THE ANDROID BUG (the KISS deep dive)
    # =========================================================================================
    s = d.slide("Defect 03 · why it only failed on one platform",
                "A Z-INDEX IS ", "NOT A Z-INDEX.",
                "The menu opened every time. It was painted underneath the video.",
                "root cause · measured on the live page · confidential")
    _rect(s, 0.55, 2.15, 6.05, 2.40, INK, LINE)
    _tb(s, 0.80, 2.35, 5.55, 0.26, "MEASURED ON THE LIVE PAGE", 9, CYAN, MONO, True)
    _tb(s, 0.80, 2.72, 5.55, 1.85,
        ".more-p    z-index 60   trapped: true\n"
        "#hd        z-index 20   position: sticky\n"
        "<video>    top: 123px\n\n"
        "sticky + z-index = a stacking context",
        10.5, BODY, MONO, space=1.45)
    bullets(s, 7.00, 2.30, 5.78, [
        "A z-index inside a stacking context orders an element only against its SIBLINGS.",
        "The panel's 60 never competed with the page — the whole header competed, at 20.",
        "Android promotes a <video> to a hardware overlay that paints over non-composited content.",
        "On /demo the video starts right below the header, so the entire panel was drawn under it.",
        "On desktop the video is position:static — identical code, visibly fine."], gap=0.50, size=10.5)
    _rect(s, 0.55, 4.95, 12.23, 1.85, PANEL, LINE)
    _tb(s, 0.85, 5.15, 5.55, 0.30, "THE FIX", 10.5, GREEN, MONO, True)
    _tb(s, 0.85, 5.50, 5.55, 1.15,
        "Render the panel through a portal to <body>: outside every ancestor stacking context "
        "and clip, on its own layer. Outside-click became a real backdrop element instead of a "
        "listener that races the opening gesture.", 10.5, BODY, TEXT, space=1.3)
    _tb(s, 6.95, 5.15, 5.55, 0.30, "THE LESSON", 10.5, AMBER, MONO, True)
    _tb(s, 6.95, 5.50, 5.55, 1.15,
        "The failure needed BOTH a stacking context and a composited overlay. Two harmless "
        "things, only fatal together — which is exactly the class of defect a single "
        "reviewer, reasoning from one mental model, is worst at finding.", 10.5, BODY, TEXT, space=1.3)

    # =========================================================================================
    # 04 — HOW THE GATE WORKS
    # =========================================================================================
    s = d.slide("The release gate · one command",
                "TEST. STAGE. ", "REBOOT. DECIDE.",
                "python ship.py — and nothing reaches production that has not survived all of it.",
                "ship.py · stagegate.py · quorum.py · confidential")
    steps = [("01", "TEST", "37 unit tests +\n14 engine suites", CYAN),
             ("02", "COMMIT", "git archive HEAD\nimmutable bytes", INDIGO),
             ("03", "STAGING TWIN", "same size, region\nand image as prod", VIOLET),
             ("04", "REBOOT IT", "the one test prod\ncan never run", AMBER),
             ("05", "35 CHECKS", "deterministic\nGO / NO-GO", GREEN),
             ("06", "4-MODEL PANEL", "written verdict\nadvisory only", CYAN)]
    for i, (num, head, body, col) in enumerate(steps):
        x = 0.55 + i * 2.08
        _rect(s, x, 2.30, 1.92, 1.72, INK, LINE)
        _rect(s, x, 2.30, 1.92, 0.05, col, None)
        _tb(s, x + 0.18, 2.48, 1.56, 0.26, num, 12, col, DISPLAY, True)
        _tb(s, x + 0.18, 2.80, 1.56, 0.30, head, 10.5, WHITE, DISPLAY, True)
        _tb(s, x + 0.18, 3.16, 1.60, 0.75, body, 8.5, BODY, TEXT, space=1.25)
        if i < 5:
            _tb(s, x + 1.94, 2.98, 0.16, 0.3, "›", 14, MUTED, TEXT)
    _rect(s, 0.55, 4.30, 12.23, 0.95, PANEL, LINE)
    _tb(s, 0.85, 4.48, 11.63, 0.30, "THE STANDING RULE", 10.5, GREEN, MONO, True)
    _tb(s, 0.85, 4.80, 11.63, 0.34,
        "The models write the reasoning. The checks decide the release. A model has never been "
        "able to wave a broken build through.", 11.5, BODY, TEXT)
    # No kicker row here: with one, the body had 0.23in for three lines and the last line sat
    # ON the card border in the QA render. Head + body only, and the card is taller.
    for i, (h, b, col) in enumerate([
            ("Fails safe both ways", "A rate-limited model cannot block a good release; an "
             "agreeable model cannot pass a dead container.", CYAN),
            ("Reboot is the point", "boot_id must CHANGE — a test that can pass without the "
             "event happening is not a test.", AMBER),
            ("Synthetic data only", "The twin is built from committed RFC 5737 fixtures. No "
             "production personal data crosses over.", VIOLET)]):
        x = 0.55 + i * 4.13
        _rect(s, x, 5.40, 3.93, 1.40, INK, LINE)
        _rect(s, x, 5.40, 3.93, 0.05, col, None)
        _tb(s, x + 0.26, 5.60, 3.41, 0.30, h, 12.5, WHITE, DISPLAY, True)
        _tb(s, x + 0.26, 5.98, 3.41, 0.72, b, 9.5, BODY, TEXT, space=1.28)

    # =========================================================================================
    # 05 — THE PANEL
    # =========================================================================================
    s = d.slide("The panel · four vendors, two roles",
                "NO SHARED ", "FAILURE DOMAIN.",
                "Two soldiers argue the release is safe. Two auditors try to break it.",
                "quorum.py · deepseek · meta · google · moonshot · confidential")
    for i, (role, model, vendor, job, col) in enumerate([
            ("SOLDIER", "deepseek-3.2", "DeepSeek", "Head of the chain. Measured fastest and "
             "contract-valid on the real 13k-character prompt.", CYAN),
            ("SOLDIER", "llama-4-maverick", "Meta", "Open weights, different lineage. Weakest "
             "reviewer of the four — documented, kept for vendor spread.", INDIGO),
            ("AUDITOR", "gemma-4-31B-it", "Google", "Twice caught a check whose own detail "
             "contradicted its verdict.", VIOLET),
            ("AUDITOR", "kimi-k2.6", "Moonshot", "Sharpest auditor on record here. Also the "
             "most prone to inventing architecture it cannot see.", AMBER)]):
        x = 0.55 + i * 3.10
        _rect(s, x, 2.20, 2.94, 2.60, INK, LINE)
        _rect(s, x, 2.20, 2.94, 0.06, col, None)
        _tb(s, x + 0.22, 2.42, 2.50, 0.24, role, 9, col, MONO, True)
        _tb(s, x + 0.22, 2.72, 2.50, 0.30, model, 11.5, WHITE, DISPLAY, True)
        _tb(s, x + 0.22, 3.08, 2.50, 0.24, vendor, 9, MUTED, MONO)
        _tb(s, x + 0.22, 3.42, 2.50, 1.25, job, 9.5, BODY, TEXT, space=1.3)
    _rect(s, 0.55, 5.00, 12.23, 1.80, PANEL, LINE)
    _tb(s, 0.85, 5.20, 11.63, 0.30, "WHY FOUR VENDORS AND NOT FOUR PROMPTS",
        10.5, CYAN, MONO, True)
    _tb(s, 0.85, 5.55, 11.63, 1.10,
        "A rate limit, an outage or a training blind spot is PROVIDER-WIDE. Four hats on one "
        "model share every one of them: ask it four times and you get four correlated answers "
        "and one correlated silence. Four vendors do not go down together, are not fine-tuned on "
        "the same data, and do not share a failure mode.\n"
        "The same rule already governs the product: the false-positive auditor is never the "
        "model that wrote the deck, and never the same vendor.", 11, BODY, TEXT, space=1.32)

    # =========================================================================================
    # 06 — SCIENCE 1
    # =========================================================================================
    s = d.slide("The evidence · 1 of 3",
                "A MODEL CANNOT ", "MARK ITS OWN HOMEWORK.",
                "The single most load-bearing result behind this design.",
                "Huang et al., ICLR 2024 · arXiv:2310.01798 · confidential")
    _rect(s, 0.55, 2.20, 6.05, 2.55, INK, LINE)
    _tb(s, 0.80, 2.42, 5.55, 0.30, "GPT-4 ON GSM8K, ASKED TO SELF-CORRECT", 9, CYAN, MONO, True)
    for i, (lab, val, col) in enumerate([("standard, 1 call", "95.5", WHITE),
                                         ("after 1 round", "91.5", AMBER),
                                         ("after 2 rounds", "89.0", RED)]):
        y = 2.85 + i * 0.58
        _tb(s, 0.80, y, 3.10, 0.32, lab, 10.5, BODY, TEXT)
        _tb(s, 4.10, y - 0.06, 1.30, 0.40, val, 17, col, DISPLAY, True, PP_ALIGN.RIGHT)
    _tb(s, 0.80, 4.36, 5.55, 0.26, "5x the compute, 6.5 points WORSE",
        10, MUTED, MONO, True)
    bullets(s, 7.00, 2.30, 5.78, [
        "Unaided self-correction does not just fail to help — it DEGRADES results.",
        "GPT-3.5 on CommonSenseQA collapses 75.8 → 38.1 the same way.",
        "The model has no independent signal, so it talks itself out of correct answers.",
        "The critical caveat, and we keep it: with EXTERNAL feedback, self-correction WORKS "
        "— GSM8K rises 95.5 → 97.5.",
        "So the value was never in the reviewing. It is in the signal being external."],
        gap=0.52, size=10.5)
    _rect(s, 0.55, 4.95, 12.23, 1.85, PANEL, LINE)
    _rect(s, 0.55, 4.95, 0.09, 1.85, CYAN, None)
    _tb(s, 0.90, 5.15, 11.55, 0.30, "WHAT WE BUILT FROM IT", 10.5, CYAN, MONO, True)
    _tb(s, 0.90, 5.50, 11.55, 1.15,
        "The panel exists to be EXTERNAL to whoever wrote the code — that is the exact "
        "condition under which the literature shows review helps. It is not a claim that four "
        "models out-think one; it is a claim that a reviewer who did not write the code is a "
        "different instrument from the author, and the paper says the difference is the whole "
        "effect.", 11, BODY, TEXT, space=1.32)

    # =========================================================================================
    # 07 — SCIENCE 2
    # =========================================================================================
    s = d.slide("The evidence · 2 of 3",
                "IT RECOGNISES ", "ITS OWN WORK.",
                "And having recognised it, it prefers it — out of proportion to its quality.",
                "Panickssery et al., NeurIPS 2024 · arXiv:2404.13076 · confidential")
    for i, (v, lab, col) in enumerate([
            ("73.5%", "GPT-4 accuracy at telling its OWN output\nfrom two other models' and a human's",
             CYAN),
            (">90%", "self-recognition after fine-tuning on\njust 500 examples", VIOLET),
            ("57–63%", "the human-rated quality gap between the\nsame models — far smaller "
             "than the self-preference", AMBER)]):
        x = 0.55 + i * 4.13
        _rect(s, x, 2.20, 3.93, 1.85, INK, LINE)
        _rect(s, x, 2.20, 3.93, 0.06, col, None)
        _tb(s, x + 0.26, 2.48, 3.41, 0.55, v, 26, col, DISPLAY, True)
        _tb(s, x + 0.26, 3.15, 3.41, 0.75, lab, 9.5, BODY, TEXT, space=1.25)
    _tb(s, 0.55, 4.25, 12.23, 0.30,
        "The authors' conclusion, closely paraphrased: the self-preference is DISPROPORTIONATE "
        "to the actual quality difference.", 11, WHITE, TEXT, True)
    _rect(s, 0.55, 4.70, 6.05, 2.10, INK, LINE)
    _tb(s, 0.80, 4.90, 5.55, 0.30, "WHAT IT MEANS FOR A RELEASE GATE", 10.5, CYAN, MONO, True)
    _tb(s, 0.80, 5.25, 5.55, 1.35,
        "An author-judge is not neutral and cannot be made neutral by prompting. So the auditor "
        "is chosen to be a different model AND a different vendor from the one that produced the "
        "work — in the gate, and in the product's own false-positive audit.",
        10.5, BODY, TEXT, space=1.32)
    _rect(s, 6.73, 4.70, 6.05, 2.10, INK, LINE)
    _tb(s, 6.98, 4.90, 5.55, 0.30, "TWO THINGS WE REFUSED TO PUT ON THIS SLIDE",
        10.5, AMBER, MONO, True)
    _tb(s, 6.98, 5.25, 5.55, 1.35,
        "• A correlation coefficient for self-recognition vs self-preference. It does not "
        "exist in the paper — the finding is shown in scatter plots.\n"
        "• The widely-repeated line that models “prefer themselves mostly "
        "legitimately”. It is not in the paper and it inverts the conclusion.",
        10, BODY, TEXT, space=1.28)

    # =========================================================================================
    # 08 — SCIENCE 3
    # =========================================================================================
    s = d.slide("The evidence · 3 of 3",
                "AND THE PART THAT ", "ARGUES AGAINST US.",
                "An honest deck includes the result that does not flatter the design.",
                "Wang et al. ICLR 2023 · Huang et al. ICLR 2024 · Khan et al. ICML 2024 · confidential")
    card(s, 0.55, 2.20, 3.95, 2.45, "supports us", GREEN, "A separate judge helps",
         "ICML 2024 Best Paper: two expert models argue opposing answers and a SEPARATE, weaker "
         "judge decides. Non-expert model judges reach 76% accuracy against a 48% baseline; "
         "human judges 88% vs 60%.\n\nKhan et al., PMLR 235.")
    card(s, 4.69, 2.20, 3.95, 2.45, "supports us", GREEN, "Diversity + voting wins",
         "Self-consistency — sample several independent reasoning paths and take the "
         "majority — lifts GSM8K by +17.9 points. Independent attempts beat one confident "
         "attempt.\n\nWang et al., ICLR 2023.")
    card(s, 8.83, 2.20, 3.95, 2.45, "argues against us", RED, "Debate loses to voting",
         "At MATCHED compute on GSM8K, multi-agent debate scores 83.0 while plain "
         "self-consistency scores 88.2. The authors attribute the gain to consistency, not to "
         "self-correction.\n\nHuang et al., ICLR 2024, Table 7.", hcol=WHITE)
    _rect(s, 0.55, 4.95, 12.23, 1.85, PANEL, LINE)
    _rect(s, 0.55, 4.95, 0.09, 1.85, AMBER, None)
    _tb(s, 0.90, 5.15, 11.55, 0.30, "SO WHAT IS ACTUALLY DEFENSIBLE", 10.5, AMBER, MONO, True)
    _tb(s, 0.90, 5.50, 11.55, 1.15,
        "Adding models is not automatically better — naive debate can lose to simple "
        "majority voting at the same cost. What the evidence supports is narrower and stronger: "
        "an INDEPENDENT judge, structurally unable to share the author's blind spot or prefer the "
        "author's work, plus a decision made by code rather than by consensus. That is precisely "
        "the shape of this gate — and it is why the models advise and never decide.",
        11, BODY, TEXT, space=1.32)

    # =========================================================================================
    # 09 — SCOREBOARD: HITS
    # =========================================================================================
    s = d.slide("The record · what the panel caught",
                "SIX CATCHES ", "I MISSED.",
                "Every one of these became a permanent check in the pipeline.",
                "verbatim from the ship logs · confidential")
    rows = [("kimi-k2.6", "Read d41d8cd98f00 as the md5 of the EMPTY STRING — and identified "
             "a disabled admin API from the hash alone.", "best single catch on record"),
            ("gemma + deepseek", "Twice spotted that a check's DETAIL contradicted its own "
             "VERDICT — the tell for “the check is broken, not the system”.",
             "both times correct"),
            ("kimi-k2.6", "“engine_runs proves invocation, not correctness.”",
             "→ became the artifact-content check"),
            ("kimi-k2.6", "“Config drift is undetected” and “one bad block takes "
             "every domain down.”", "→ became config_drift + a negative test"),
            ("kimi-k2.6", "Nothing asserts WHICH domains are served.",
             "→ became agent.py roster"),
            ("kimi + gemma", "Flagged vhost_roster SKIPPED in today's run — a check that "
             "reports SKIP is not a check.", "open, and correctly raised")]
    for i, (who, what, outcome) in enumerate(rows):
        y = 2.18 + i * 0.735
        _rect(s, 0.55, y, 12.23, 0.66, INK if i % 2 == 0 else PANEL, LINE)
        _rect(s, 0.55, y, 0.06, 0.66, GREEN, None)
        _tb(s, 0.80, y + 0.09, 1.85, 0.24, who, 9, CYAN, MONO, True)
        _tb(s, 0.80, y + 0.32, 8.60, 0.30, what, 10, BODY, TEXT)
        _tb(s, 9.55, y + 0.22, 3.10, 0.36, outcome, 9, GREEN, MONO, space=1.2)
    # Rows end at 6.59. This line used to sit at 6.95 and printed straight through the footer
    # at 7.04 — two runs interleaved, both unreadable. Footers are load-bearing furniture.
    _tb(s, 0.55, 6.70, 12.23, 0.28,
        "Standing rule: when the panel is right, it becomes a check in the SAME change — "
        "not a note, not a “known gap”.", 10, MUTED, TEXT)

    # =========================================================================================
    # 10 — SCOREBOARD: MISSES
    # =========================================================================================
    s = d.slide("The record · what the panel got wrong",
                "AND WHERE IT ", "IS UNRELIABLE.",
                "The same panel, the same runs. This is why it advises and does not decide.",
                "verbatim from the ship logs · confidential")
    rows = [("kimi-k2.6", "Inverted config_reread — argued “started AFTER the write "
             "means stale”, the exact opposite of the truth.", "refuted by its own evidence"),
            ("kimi-k2.6", "Invented an engine job queue, an ENGINE_MODE variable and a "
             "Kubernetes manifest.", "none of them exist here"),
            ("llama-4-maverick", "Repeatedly restates the failure as its own diagnosis "
             "(“the engine is not running correctly”).", "adds no information"),
            ("ALL FOUR", "On the config_drift false positive, every model proposed a confident "
             "fix for a fault that did not exist.", "the check was broken, not the box"),
            ("kimi-k2.6", "Today: called config_drift “suspect by construction” — "
             "reasoning from a hash comparison removed hours earlier.", "stale premise")]
    for i, (who, what, outcome) in enumerate(rows):
        y = 2.18 + i * 0.82
        _rect(s, 0.55, y, 12.23, 0.72, INK if i % 2 == 0 else PANEL, LINE)
        _rect(s, 0.55, y, 0.06, 0.72, RED, None)
        _tb(s, 0.80, y + 0.11, 1.85, 0.24, who, 9, AMBER, MONO, True)
        _tb(s, 0.80, y + 0.35, 8.60, 0.30, what, 10, BODY, TEXT)
        _tb(s, 9.55, y + 0.25, 3.10, 0.36, outcome, 9, RED, MONO, space=1.2)
    _rect(s, 0.55, 6.35, 12.23, 0.58, PANEL, LINE)
    _tb(s, 0.85, 6.46, 11.63, 0.36,
        "The pattern is consistent: strong reasoning FROM evidence in front of it, unreliable "
        "when extrapolating to architecture it cannot see. So: feed it more evidence, never more "
        "authority.", 10.5, WHITE, TEXT)

    # =========================================================================================
    # 11 — THE DAY THE PANEL WAS RIGHT
    # =========================================================================================
    s = d.slide("The strongest evidence we have",
                "4 OF 4 SAID NO-GO. ", "IT SHIPPED ANYWAY.",
                "The panel was right and the deterministic gate was lying.",
                "incident 7 Aug 2026 · now guarded by tests/test_gate_integrity.py · confidential")
    _rect(s, 0.55, 2.20, 12.23, 1.42, INK, LINE)
    _tb(s, 0.80, 2.40, 11.73, 0.26, "THE TWO LINES FROM THE RUN", 9, CYAN, MONO, True)
    _tb(s, 0.80, 2.72, 11.73, 0.80,
        "mount_fresh    OK   container reads the current file (3af610cdeb71) - bind mount is not stale\n"
        "config_drift   OK   drift check unavailable: STALE MOUNT STALE MOUNT: host= container=3af610cdeb71",
        9.5, BODY, MONO, space=1.45)
    for i, (n, h, b, col) in enumerate([
            ("01", "A hard-coded path", "The agent hashed production's file path, so on staging "
             "it read nothing — and reported “no file here” as “the mount is "
             "stale”.", RED),
            ("02", "A catch-all that PASSED", "An unrecognised verdict fell through to a branch "
             "that scored it as a pass. A fallback that turns an unknown answer into a success is "
             "worse than no check.", RED),
            ("03", "Nobody read the detail", "A check reporting PASS while its own text said "
             "“STALE MOUNT”. All four models saw it. The code did not.", AMBER)]):
        x = 0.55 + i * 4.13
        card(s, x, 3.78, 3.93, 1.75, ("defect " + n), col, h, b, bsize=9.5)
    _rect(s, 0.55, 5.72, 12.23, 1.08, PANEL, LINE)
    _rect(s, 0.55, 5.72, 0.09, 1.08, GREEN, None)
    _tb(s, 0.90, 5.86, 11.55, 0.30, "WHAT CHANGED — NARROWLY, AND ON PURPOSE",
        10.5, GREEN, MONO, True)
    _tb(s, 0.90, 6.18, 11.55, 0.52,
        "Models still cannot veto a release. But a UNANIMOUS panel against a GREEN gate now "
        "HALTS the ship and requires a human override — because twice now that exact "
        "pattern has meant a check is lying. A split panel still does not block.",
        10.5, BODY, TEXT, space=1.28)

    # =========================================================================================
    # 12 — WHAT WE DO NOT CLAIM
    # =========================================================================================
    s = d.slide("Honesty · the boundary of the claim",
                "WHAT WE DO ", "NOT CLAIM.",
                "The product's credibility rests on “absence of evidence is never a "
                "finding”. The marketing has to obey the same rule.",
                "confidential")
    card(s, 0.55, 2.20, 6.05, 2.15, "not claimed", RED,
         "“It beats Claude”",
         "There is no benchmark in this repo comparing the panel to any named product. The panel "
         "reviews a running system against deterministic evidence; a coding model writes the "
         "code. They are not measured on the same task, so the comparison has no number behind "
         "it — and an unmeasured superiority claim against a named competitor also needs "
         "substantiation under UWG §6 and the UCP Directive.")
    card(s, 6.73, 2.20, 6.05, 2.15, "claimed, and evidenced", GREEN,
         "It catches what one reviewer misses",
         "Six documented catches that the author (me) missed, each now a permanent check. One "
         "documented case where the panel was unanimously right and the deterministic gate was "
         "wrong. Five documented cases where the panel was wrong. All of it is in the repo, "
         "dated, in the operator's own logs.")
    _rect(s, 0.55, 4.55, 12.23, 2.25, PANEL, LINE)
    _tb(s, 0.85, 4.75, 11.63, 0.30, "THE MECHANISM, IN ONE SENTENCE", 10.5, CYAN, MONO, True)
    _tb(s, 0.85, 5.10, 11.63, 0.60,
        "A reviewer that did not write the code, cannot recognise it as its own, and does not "
        "share the author's vendor, training data or outage window is a genuinely INDEPENDENT "
        "instrument — and independence, not intelligence, is what the evidence says makes "
        "review work.", 12.5, WHITE, TEXT, True, space=1.3)
    _tb(s, 0.85, 5.90, 11.63, 0.80,
        "If a head-to-head comparison is wanted later, the honest route is compare_models.py "
        "extended to a published methodology and a dated result. That would be worth far more "
        "than the assertion — and it is the only version of this claim we would put in "
        "front of a bank.", 10.5, BODY, TEXT, space=1.3)

    # =========================================================================================
    # 13 — SOURCES
    # =========================================================================================
    s = d.slide("Sources · every claim on slides 06–08",
                "CHECK IT ", "YOURSELF.",
                "Verified against the primary papers, not against secondary summaries.",
                "confidential")
    src = [("Huang et al. (2024)", "Large Language Models Cannot Self-Correct Reasoning Yet",
            "ICLR 2024 · arXiv:2310.01798", "GPT-4 GSM8K 95.5 → 89.0 unaided (Table 3); "
            "95.5 → 97.5 with external feedback (Table 2); debate 83.0 vs voting 88.2 (Table 7)"),
           ("Panickssery et al. (2024)", "LLM Evaluators Recognize and Favor Their Own Generations",
            "NeurIPS 2024 · arXiv:2404.13076", "GPT-4 self-recognition 73.5%; >90% after 500 "
            "fine-tuning examples; self-preference disproportionate to human-rated quality"),
           ("Wang et al. (2023)", "Self-Consistency Improves Chain of Thought Reasoning",
            "ICLR 2023 · arXiv:2203.11171", "GSM8K +17.9, SVAMP +11.0, AQuA +12.2"),
           ("Khan et al. (2024)", "Debating with More Persuasive LLMs Leads to More Truthful Answers",
            "ICML 2024 Best Paper · PMLR 235", "Non-expert model judges 76% vs 48% baseline; "
            "human judges 88% vs 60%"),
           ("Zheng et al. (2023)", "Judging LLM-as-a-Judge with MT-Bench and Chatbot Arena",
            "NeurIPS 2023 D&B · arXiv:2306.05685", "Named self-enhancement bias and reported "
            "it INCONCLUSIVE — cited here for the term, not as proof")]
    for i, (who, title, venue, finding) in enumerate(src):
        y = 2.18 + i * 0.86
        _rect(s, 0.55, y, 12.23, 0.78, INK if i % 2 == 0 else PANEL, LINE)
        _tb(s, 0.80, y + 0.11, 2.45, 0.24, who, 9.5, CYAN, MONO, True)
        _tb(s, 0.80, y + 0.38, 6.20, 0.30, title, 10, WHITE, TEXT, True)
        _tb(s, 7.15, y + 0.11, 5.40, 0.24, venue, 8.5, MUTED, MONO)
        _tb(s, 7.15, y + 0.38, 5.40, 0.32, finding, 8.5, BODY, TEXT, space=1.2)
    _rect(s, 0.55, 6.55, 12.23, 0.38, PANEL, LINE)
    _tb(s, 0.85, 6.63, 11.63, 0.24,
        "Two claims were dropped in fact-checking: a correlation coefficient that does not exist "
        "in Panickssery, and a cross-examination paper that does not test cross-MODEL review.",
        9, AMBER, TEXT)

    return d.save(out)


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--template", default=os.path.join(
        os.path.dirname(os.path.abspath(__file__)),
        "S4biz_Sovereign_Cyber_Cloud_Capability_Brief.pptx"))
    ap.add_argument("--out", default="S4biz_Consensus_Deep_Dive.pptx")
    a = ap.parse_args()
    print("built:", build(a.template, a.out))
