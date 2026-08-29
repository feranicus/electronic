#!/usr/bin/env python3
"""deck_chrome.py — the guards and geometry every S4biz deck shares. One implementation.

WHY THIS FILE EXISTS
    `_check_title` had been copied into four builders and `_tail_cyan` into two. Every copy is a
    place the rule can drift, and this repository has paid for that pattern more than any other:
    ENRICH_MODELS had four homes, the document language set had six, and the container-to-host
    path resolution was fixed in dbbackup and then got wrong again in logship because the fix was
    not reachable from there.

    The TEMPLATE itself (Deck, card, bullets, stat, the palette, the wordmark chrome) stays in
    build_consensus_deck.py, which read it off the S4biz capability brief. This module holds only
    the things layered ON TOP of it by the decks that are proposals rather than briefs.

WHAT IS IN HERE, AND WHY EACH ONE IS A REFUSAL RATHER THAN A CONVENTION

    check_title   A fixed-height title row is arithmetic. The box is 0.70in at 30pt, so a second
                  line needs about 1.04in and lands on the sub-heading at y=1.55. Measured on the
                  render: 49 characters fit, 53 wrap. The cap is 50, the lower end of that
                  interval, chosen because it is observed to fit and not because it lets the next
                  title through.

    tail_colour   Recolours the title tail without forking the shared Deck implementation.
                  IT RETURNS WHETHER IT DID ANYTHING, and callers are expected to check. The first
                  version of this in the legal deck looked for a paragraph with exactly two runs.
                  `_tb` adds a run even for an empty string, so the title paragraph actually has
                  three, the function matched nothing, every tail stayed violet, and the build
                  printed "built" and exited zero. A cosmetic pass that cannot fail is not a pass.

    geometry      Column arithmetic computed once. Three slides inventing their own column widths
                  is three chances to be 0.02in out and only notice it in a render.
"""
from build_consensus_deck import CYAN  # noqa: F401  — re-exported as the default tail colour
from pptx.util import Pt

TITLE_MAX = 50

# Column geometry. Left margin 0.55, right edge 12.78, gaps chosen so each row ends flush.
X2, W2 = (0.55, 6.78), 6.00
X3, W3 = (0.55, 4.69, 8.83), 3.90
X4, W4 = (0.55, 3.64, 6.73, 9.82), 2.87


def check_title(title, tail):
    n = len(title) + len(tail or "")
    if n > TITLE_MAX:
        raise SystemExit("[X] title is %d characters and wraps onto the sub-heading (max %d): %r"
                         % (n, TITLE_MAX, (title + (tail or ""))))
    return n


def tail_colour(slide, colour):
    """Recolour the title tail. Returns True if a tail was found and changed.

    The title runs are the only 30pt runs on a non-hero slide, so the box is identified by
    measurement rather than by index, and the LAST such run is the tail.
    """
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            big = [r for r in p.runs if r.font.size == Pt(30) and r.text.strip()]
            if len(big) >= 2:
                big[-1].font.color.rgb = colour
                return True
    return False


FOOTER_Y = 7.04          # where Deck.slide() puts the footer and the page number
_EPS = 0.005             # float tolerance when deciding "this IS the footer"


def assert_above_footer(slide, label):
    """Refuse any body text that runs into the footer band.

    THE DEFECT THIS CATCHES, which this repository has now paid for four separate times: a fixed
    row height times a row count is arithmetic, and getting it wrong puts text through the footer.
    It happened on the site header row twice, on the roadmap deck's six tier rows, and again on
    the four-page proposal, where six price rows at a 0.52in step ended at 7.01 against a footer
    at 7.04. Eyeballing a render catches it only if somebody looks at that corner.

    Only text is checked. The background picture legitimately covers the whole slide, and the
    footer and page number legitimately sit in the band.

    THE FIRST VERSION OF THIS GUARD HAD A HOLE THAT WOULD HAVE MISSED THE VERY BUG IT WAS WRITTEN
    FOR. It skipped any shape starting at or after 7.00, and the offending note sat at exactly
    7.00 with a 0.28in height, so it was excused rather than flagged. The exemption must identify
    the FOOTER, which starts at 7.04, and nothing else. Found by a negative test that did not
    fail, which is the only reason to write negative tests at all.
    """
    bad = []
    for sh in slide.shapes:
        # PICTURES ARE EXEMPT and nothing else is. The background image legitimately covers the
        # whole slide. A CARD is not exempt: a panel sliding under the footer puts the footer text
        # on top of a filled rectangle, which is just as visible as overflowing text. The first
        # version checked text only, and two of three negative tests walked straight past it.
        if sh.shape_type is not None and str(sh.shape_type).startswith("PICTURE"):
            continue
        txt = sh.text_frame.text.strip() if sh.has_text_frame else ""
        top = sh.top / 914400.0
        bottom = top + sh.height / 914400.0
        if top >= FOOTER_Y - _EPS:       # this shape IS the footer or the page number
            continue
        if bottom > FOOTER_Y:
            bad.append((round(top, 2), round(bottom, 2),
                        txt[:55].replace("\n", " ") or "<panel or rule, no text>"))
    if bad:
        raise SystemExit(
            "[X] %r: %d text box(es) run into the footer at %.2fin. Reduce the row step or the "
            "block height; do not nudge the footer.\n    %s"
            % (label, len(bad), FOOTER_Y,
               "\n    ".join("y %.2f..%.2f  %r" % b for b in bad)))


def assert_layout(deck):
    """Run the footer check over every slide. Call it immediately before save().

    It cannot run inside the slide() wrapper, because at that point the slide holds only the
    chrome and none of the content that actually overflows.
    """
    for i, s in enumerate(deck.prs.slides, 1):
        assert_above_footer(s, "slide %02d" % i)
    return deck


def guarded(deck, tail=None):
    """Install the title guard, and optionally the tail recolour, on a Deck.

    Wrapping the deck's own method means a NEW slide cannot forget to call the guard, which is the
    difference between a rule and a habit.
    """
    original = deck.slide

    def _slide(eyebrow, title, title_tail=None, sub=None, footer="", hero=False):
        if not hero:
            check_title(title, title_tail)
        s = original(eyebrow, title, title_tail, sub, footer, hero)
        if tail is not None and not hero and title_tail and not tail_colour(s, tail):
            raise SystemExit(
                "[X] the title tail was not recoloured on %r. Deck.slide()'s run layout has "
                "changed and this deck's colour pass is now a silent no-op." % title)
        return s

    deck.slide = _slide
    return deck
